from pathlib import Path
import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("/Users/inti/GitHub/clases/modulo_1/clase_7_apoyo_docente_llm.pptx")

PAPER = RGBColor(249, 245, 236)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(35, 38, 45)
SLATE = RGBColor(96, 102, 117)
CLAY = RGBColor(183, 93, 63)
CLAY_SOFT = RGBColor(247, 232, 226)
TEAL = RGBColor(35, 120, 115)
TEAL_SOFT = RGBColor(224, 245, 242)
GOLD = RGBColor(206, 156, 49)
GOLD_SOFT = RGBColor(251, 243, 214)
BLUE = RGBColor(74, 102, 173)
BLUE_SOFT = RGBColor(228, 236, 250)
GREEN = RGBColor(66, 145, 98)
GREEN_SOFT = RGBColor(229, 244, 235)
ROSE = RGBColor(191, 92, 87)
ROSE_SOFT = RGBColor(249, 232, 230)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Aptos",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    lines = text if isinstance(text, list) else str(text).split("\n")
    for index, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.name = font_name
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(2)
    return box


def add_decor(slide, color=CLAY):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.18),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()


def add_title(slide, title, subtitle=None, tag=None):
    if tag:
        add_tag(slide, tag, 0.72, 0.36, fill_rgb=CLAY_SOFT, text_rgb=CLAY, min_width=1.45)
    add_text(slide, title, 0.72, 0.72, 11.2, 0.7, font_size=28, bold=True, font_name="Aptos Display")
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.28, 11.3, 0.45, font_size=12.5, color=SLATE)


def add_footer(slide, text):
    add_text(slide, text, 0.72, 7.02, 11.8, 0.2, font_size=9, color=SLATE, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, fill_rgb=WHITE, line_rgb=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb or fill_rgb
    shape.line.width = Pt(1.2)
    return shape


def add_tag(slide, text, left, top, fill_rgb=TEAL_SOFT, text_rgb=TEAL, min_width=0.95):
    width = max(min_width, 0.12 * len(text) + 0.35)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.34),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = fill_rgb
    add_text(slide, text, left, top + 0.02, width, 0.2, font_size=10, color=text_rgb, bold=True, align=PP_ALIGN.CENTER)
    return width


def estimate_chip_width(token):
    return max(0.42, min(1.6, 0.12 * len(token) + 0.34))


def add_token_cluster(
    slide,
    tokens,
    left,
    top,
    max_width,
    fill_rgb=BLUE_SOFT,
    text_rgb=INK,
    highlight=None,
    highlight_fill=None,
    chip_height=0.38,
):
    x_pos = left
    y_pos = top
    highlight = set(highlight or [])
    highlight_fill = highlight_fill or GOLD_SOFT
    for index, token in enumerate(tokens):
        chip_width = estimate_chip_width(token)
        if x_pos + chip_width > left + max_width:
            x_pos = left
            y_pos += chip_height + 0.12
        chip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x_pos),
            Inches(y_pos),
            Inches(chip_width),
            Inches(chip_height),
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = highlight_fill if index in highlight else fill_rgb
        chip.line.color.rgb = chip.fill.fore_color.rgb
        add_text(
            slide,
            token,
            x_pos,
            y_pos + 0.02,
            chip_width,
            chip_height,
            font_size=11,
            color=text_rgb,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        x_pos += chip_width + 0.08
    return y_pos + chip_height


def add_candidate(slide, token, probability, left, top, fill_rgb, accent_rgb):
    add_card(slide, left, top, 2.82, 1.7, fill_rgb=fill_rgb, line_rgb=accent_rgb)
    add_text(slide, token, left + 0.18, top + 0.22, 2.46, 0.45, font_size=19, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, f"{int(probability * 100)}%", left + 0.18, top + 0.7, 2.46, 0.28, font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
    track = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.28),
        Inches(top + 1.15),
        Inches(2.26),
        Inches(0.18),
    )
    track.fill.solid()
    track.fill.fore_color.rgb = PAPER
    track.line.color.rgb = PAPER
    value = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 0.28),
        Inches(top + 1.15),
        Inches(2.26 * probability),
        Inches(0.18),
    )
    value.fill.solid()
    value.fill.fore_color.rgb = accent_rgb
    value.line.color.rgb = accent_rgb


def add_bar_group(slide, title, subtitle, values, left, top):
    add_card(slide, left, top, 3.55, 3.3, fill_rgb=WHITE, line_rgb=BLUE_SOFT)
    add_text(slide, title, left + 0.2, top + 0.18, 3.1, 0.28, font_size=15, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, subtitle, left + 0.2, top + 0.48, 3.1, 0.25, font_size=10.5, color=SLATE, align=PP_ALIGN.CENTER)
    base_y = top + 2.65
    chart_height = 1.5
    max_value = max(value for _, value, _ in values)
    for index, (label, value, color) in enumerate(values):
        bar_x = left + 0.45 + index * 0.95
        bar_height = chart_height * (value / max_value)
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(bar_x),
            Inches(base_y - bar_height),
            Inches(0.52),
            Inches(bar_height),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.color.rgb = color
        add_text(slide, f"{int(round(value * 100))}%", bar_x - 0.05, base_y - bar_height - 0.28, 0.62, 0.2, font_size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, bar_x - 0.12, base_y + 0.06, 0.75, 0.35, font_size=10, align=PP_ALIGN.CENTER)


def add_memory_bar(slide, label, value, max_value, left, top, width=4.5, color=TEAL):
    add_text(slide, label, left, top - 0.03, 1.3, 0.22, font_size=12, bold=True)
    track = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 1.35),
        Inches(top),
        Inches(width),
        Inches(0.22),
    )
    track.fill.solid()
    track.fill.fore_color.rgb = PAPER
    track.line.color.rgb = PAPER
    filled = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left + 1.35),
        Inches(top),
        Inches(width * (value / max_value)),
        Inches(0.22),
    )
    filled.fill.solid()
    filled.fill.fore_color.rgb = color
    filled.line.color.rgb = color
    add_text(slide, f"{value:.1f} GB", left + 5.98, top - 0.02, 0.75, 0.22, font_size=11, color=SLATE, align=PP_ALIGN.RIGHT)


def add_flow_box(slide, label, sublabel, left, top, width=1.72, height=1.05, fill_rgb=WHITE, line_rgb=CLAY):
    add_card(slide, left, top, width, height, fill_rgb=fill_rgb, line_rgb=line_rgb)
    add_text(slide, label, left + 0.14, top + 0.16, width - 0.28, 0.28, font_size=14, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, sublabel, left + 0.14, top + 0.48, width - 0.28, 0.32, font_size=10.5, color=SLATE, align=PP_ALIGN.CENTER)


def add_timeline(slide, steps, left, top, width):
    base = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top + 0.22),
        Inches(width),
        Inches(0.05),
    )
    base.fill.solid()
    base.fill.fore_color.rgb = BLUE_SOFT
    base.line.color.rgb = BLUE_SOFT
    gap = width / (len(steps) - 1)
    for index, (time_label, title, tone) in enumerate(steps):
        x_pos = left + index * gap
        node = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x_pos - 0.18),
            Inches(top),
            Inches(0.44),
            Inches(0.44),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = tone
        node.line.color.rgb = tone
        add_text(slide, time_label, x_pos - 0.42, top + 0.55, 0.85, 0.22, font_size=10, color=tone, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x_pos - 0.63, top + 0.78, 1.3, 0.4, font_size=10.5, align=PP_ALIGN.CENTER)


def softmax(values, temperature):
    scaled = {token: score / temperature for token, score in values.items()}
    max_score = max(scaled.values())
    exp_values = {token: math.exp(score - max_score) for token, score in scaled.items()}
    total = sum(exp_values.values())
    return {token: score / total for token, score in exp_values.items()}


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    slide = prs.slides.add_slide(layout)
    set_bg(slide, PAPER)
    add_decor(slide, CLAY)
    add_title(
        slide,
        "Clase 7 | LLMs sin magia",
        "Apoyo docente rediseñado para explicar con escenas, comparaciones y flujo visual.",
        tag="APOYO DOCENTE",
    )
    loop_titles = [
        ("CONTEXTO", "lo que el modelo\nve en ese paso", CLAY_SOFT, CLAY),
        ("PROXIMO TOKEN", "elige una opcion\npor probabilidad", TEAL_SOFT, TEAL),
        ("NUEVO CONTEXTO", "la salida vuelve a\nalimentar el ciclo", GOLD_SOFT, GOLD),
    ]
    x_positions = [0.95, 4.65, 8.35]
    for x_pos, (title, subtitle, fill_rgb, tone) in zip(x_positions, loop_titles):
        add_card(slide, x_pos, 2.15, 3.1, 2.1, fill_rgb=fill_rgb, line_rgb=tone)
        add_text(slide, title, x_pos + 0.2, 2.48, 2.7, 0.42, font_size=21, bold=True, color=tone, align=PP_ALIGN.CENTER)
        add_text(slide, subtitle, x_pos + 0.25, 3.0, 2.6, 0.6, font_size=13, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, "->", 4.15, 2.82, 0.35, 0.3, font_size=28, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "->", 7.85, 2.82, 0.35, 0.3, font_size=28, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 1.05, 4.8, 11.1, 1.25, fill_rgb=WHITE, line_rgb=BLUE_SOFT)
    add_text(
        slide,
        "Idea-fuerza para decir en voz alta: un LLM no responde por magia. Repite una pregunta muy simple muchas veces: que viene despues?",
        1.35,
        5.1,
        10.5,
        0.45,
        font_size=17,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Abrir la clase señalando este loop y volver a el durante toda la hora")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, WHITE)
    add_decor(slide)
    add_title(
        slide,
        "La imagen correcta: un autocomplete probabilistico",
        "Una escena simple para reemplazar la idea de caja magica.",
        tag="INTUICION",
    )
    add_card(slide, 0.9, 1.95, 11.45, 1.05, fill_rgb=CLAY_SOFT, line_rgb=CLAY)
    add_text(slide, "el mate se toma con ...", 1.15, 2.22, 10.95, 0.4, font_size=26, bold=True, color=CLAY, align=PP_ALIGN.CENTER)
    candidates = [
        ("agua caliente", 0.55, TEAL_SOFT, TEAL),
        ("bombilla", 0.25, BLUE_SOFT, BLUE),
        ("amigos", 0.10, GOLD_SOFT, GOLD),
        ("otras", 0.10, ROSE_SOFT, ROSE),
    ]
    for index, (token, probability, fill_rgb, accent_rgb) in enumerate(candidates):
        add_candidate(slide, token, probability, 0.95 + index * 2.98, 3.4, fill_rgb, accent_rgb)
    add_card(slide, 0.95, 5.65, 11.35, 0.85, fill_rgb=PAPER, line_rgb=PAPER)
    add_text(slide, "No adivina ni entiende como una persona: ordena continuaciones posibles y toma una.", 1.2, 5.92, 10.85, 0.25, font_size=16, align=PP_ALIGN.CENTER)
    add_footer(slide, "Frase util: la sorpresa viene de la escala, no de una tarea base misteriosa")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, PAPER)
    add_decor(slide, TEAL)
    add_title(
        slide,
        "Tokenizacion: la misma frase, tres recortes distintos",
        "Mostrar esto antes de definir vocabulario: primero que vean el cambio de piezas.",
        tag="TOKENIZACION",
    )
    cards = [
        (0.75, "Character-level", "muchos pasos", list("inteligencia artificial"), BLUE_SOFT, BLUE),
        (4.45, "Word-level", "2 tokens", ["inteligencia", "artificial"], GOLD_SOFT, GOLD),
        (8.15, "Subword", "4 tokens", ["intelig", "encia", "art", "ificial"], TEAL_SOFT, TEAL),
    ]
    for left, title, badge, tokens, fill_rgb, tone in cards:
        add_card(slide, left, 1.95, 3.45, 4.15, fill_rgb=WHITE, line_rgb=fill_rgb)
        add_text(slide, title, left + 0.18, 2.15, 2.2, 0.25, font_size=17, bold=True, color=tone)
        add_tag(slide, badge, left + 2.45, 2.12, fill_rgb=fill_rgb, text_rgb=tone, min_width=0.9)
        bottom = add_token_cluster(slide, tokens, left + 0.18, 2.7, 3.05, fill_rgb=fill_rgb, highlight_fill=fill_rgb)
        add_text(slide, f"Cantidad: {len(tokens)}", left + 0.22, bottom + 0.25, 1.8, 0.2, font_size=12, color=SLATE, bold=True)
    add_footer(slide, "Cierre sugerido: token no siempre significa palabra")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, WHITE)
    add_decor(slide, GOLD)
    add_title(
        slide,
        "Contexto: lo que si ve cambia mucho la prediccion",
        "Conviene contrastar un contexto rico contra uno casi vacio.",
        tag="CONTEXTO",
    )
    add_card(slide, 0.88, 1.95, 7.3, 3.5, fill_rgb=GREEN_SOFT, line_rgb=GREEN)
    add_text(slide, "Contexto rico", 1.15, 2.18, 2.2, 0.25, font_size=17, bold=True, color=GREEN)
    add_token_cluster(
        slide,
        ["la", "capital", "de", "Argentina", "es"],
        1.18,
        2.75,
        4.9,
        fill_rgb=WHITE,
        highlight=[0, 1, 2, 3, 4],
        highlight_fill=GOLD_SOFT,
    )
    add_text(slide, "->", 6.2, 2.92, 0.35, 0.25, font_size=28, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_token_cluster(slide, ["Buenos", "Aires"], 6.58, 2.76, 1.3, fill_rgb=TEAL_SOFT, text_rgb=TEAL, highlight=[0, 1], highlight_fill=TEAL_SOFT)
    add_text(slide, "Con suficiente pista, la opcion fuerte aparece rapido.", 1.18, 4.45, 5.6, 0.25, font_size=14, color=INK)
    add_card(slide, 8.45, 1.95, 3.9, 3.5, fill_rgb=ROSE_SOFT, line_rgb=ROSE)
    add_text(slide, "Contexto pobre", 8.72, 2.18, 2.2, 0.25, font_size=17, bold=True, color=ROSE)
    add_token_cluster(slide, ["es"], 8.76, 2.8, 0.7, fill_rgb=WHITE, highlight=[0], highlight_fill=WHITE)
    add_text(slide, "->", 9.55, 2.95, 0.35, 0.25, font_size=28, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_token_cluster(slide, ["algo", "casi", "cualquiera"], 9.95, 2.78, 2.0, fill_rgb=WHITE)
    add_text(slide, "Si la ventana ve poco, la incertidumbre sube mucho.", 8.76, 4.45, 2.9, 0.35, font_size=14, color=INK)
    add_footer(slide, "Usar esta diapositiva para explicar por que la ventana de contexto importa")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, PAPER)
    add_decor(slide, BLUE)
    add_title(
        slide,
        "Probabilidad y temperatura",
        "Misma informacion de base. Distinto modo de repartir la salida.",
        tag="PROBABILIDAD",
    )
    logits = {"agua": 3.2, "yerba": 2.4, "hielo": 1.2}
    palettes = {"agua": TEAL, "yerba": GOLD, "hielo": CLAY}
    temps = [(0.3, "temperatura baja", "mas cerrada"), (1.0, "temperatura media", "equilibrio"), (1.8, "temperatura alta", "mas abierta")]
    for index, (temp, title, subtitle) in enumerate(temps):
        distribution = softmax(logits, temp)
        entries = [(token, distribution[token], palettes[token]) for token in ["agua", "yerba", "hielo"]]
        add_bar_group(slide, title, subtitle, entries, 0.9 + index * 3.95, 2.05)
    add_card(slide, 2.2, 5.7, 8.95, 0.72, fill_rgb=WHITE, line_rgb=WHITE)
    add_text(slide, "Mensaje docente: la temperatura no agrega conocimiento nuevo. Solo cambia como se usa la distribucion.", 2.45, 5.93, 8.45, 0.25, font_size=15, align=PP_ALIGN.CENTER)
    add_footer(slide, "Conectar esta idea con la clase practica de prompts y sampling")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, WHITE)
    add_decor(slide, CLAY)
    add_title(
        slide,
        "Transformer: un mapa de recorrido, no una formula",
        "La meta aca es que se lleven un flujo mental claro.",
        tag="TRANSFORMER",
    )
    steps = [
        ("Texto", "frase cruda", CLAY_SOFT, CLAY),
        ("Tokens", "piezas", BLUE_SOFT, BLUE),
        ("Embeddings", "vectores", GOLD_SOFT, GOLD),
        ("Attention", "que importa mas", TEAL_SOFT, TEAL),
        ("Feed-forward", "ajuste interno", CLAY_SOFT, CLAY),
        ("Salida", "probabilidades", GREEN_SOFT, GREEN),
    ]
    for index, (label, sublabel, fill_rgb, tone) in enumerate(steps):
        left = 0.65 + index * 2.1
        add_flow_box(slide, label, sublabel, left, 2.3, fill_rgb=fill_rgb, line_rgb=tone)
        if index < len(steps) - 1:
            add_text(slide, "->", left + 1.78, 2.66, 0.28, 0.2, font_size=24, color=SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 0.92, 4.55, 3.65, 1.2, fill_rgb=GOLD_SOFT, line_rgb=GOLD)
    add_text(slide, "Posicion", 1.14, 4.76, 1.2, 0.22, font_size=15, bold=True, color=GOLD)
    add_text(slide, "Sin orden, el modelo veria piezas pero no secuencia.", 1.14, 5.07, 3.0, 0.3, font_size=12.5)
    add_card(slide, 4.85, 4.55, 3.65, 1.2, fill_rgb=TEAL_SOFT, line_rgb=TEAL)
    add_text(slide, "Attention", 5.07, 4.76, 1.2, 0.22, font_size=15, bold=True, color=TEAL)
    add_text(slide, "No todos los tokens previos pesan igual en cada paso.", 5.07, 5.07, 3.0, 0.3, font_size=12.5)
    add_card(slide, 8.78, 4.55, 3.65, 1.2, fill_rgb=CLAY_SOFT, line_rgb=CLAY)
    add_text(slide, "Que no hace falta hoy", 9.0, 4.76, 2.1, 0.22, font_size=15, bold=True, color=CLAY)
    add_text(slide, "No hundirse en Q, K y V si todavia no fijaron la intuicion.", 9.0, 5.07, 3.0, 0.3, font_size=12.5)
    add_footer(slide, "Si entienden este recorrido, la formula puede venir despues")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, PAPER)
    add_decor(slide, TEAL)
    add_title(
        slide,
        "Cuantizacion y llama.cpp: por que esto baja a una compu normal",
        "La comparacion visual ayuda mucho mas que una definicion larga.",
        tag="LOCAL",
    )
    add_card(slide, 0.88, 1.95, 7.0, 3.7, fill_rgb=WHITE, line_rgb=BLUE_SOFT)
    add_text(slide, "Mismo modelo 7B, distinta huella de memoria", 1.15, 2.18, 3.8, 0.25, font_size=17, bold=True)
    add_memory_bar(slide, "FP16", 13.0, 13.0, 1.2, 2.9, color=ROSE)
    add_memory_bar(slide, "INT8", 6.5, 13.0, 1.2, 3.45, color=BLUE)
    add_memory_bar(slide, "INT4", 3.3, 13.0, 1.2, 4.0, color=TEAL)
    add_text(slide, "Cuantizar = usar menos bits por parametro para que el modelo pese menos y pida menos RAM.", 1.2, 4.75, 5.95, 0.45, font_size=13.5)
    add_card(slide, 8.2, 1.95, 4.15, 1.62, fill_rgb=ROSE_SOFT, line_rgb=ROSE)
    add_text(slide, "API remota", 8.48, 2.16, 1.2, 0.22, font_size=16, bold=True, color=ROSE)
    add_tag(slide, "costo variable", 8.5, 2.56, fill_rgb=WHITE, text_rgb=INK, min_width=1.15)
    add_tag(slide, "menos control", 9.95, 2.56, fill_rgb=WHITE, text_rgb=INK, min_width=1.1)
    add_tag(slide, "arranque rapido", 8.5, 2.98, fill_rgb=WHITE, text_rgb=INK, min_width=1.22)
    add_card(slide, 8.2, 4.03, 4.15, 1.62, fill_rgb=TEAL_SOFT, line_rgb=TEAL)
    add_text(slide, "llama.cpp local", 8.48, 4.24, 1.65, 0.22, font_size=16, bold=True, color=TEAL)
    add_tag(slide, "mas privacidad", 8.5, 4.64, fill_rgb=WHITE, text_rgb=INK, min_width=1.18)
    add_tag(slide, "GGUF liviano", 9.95, 4.64, fill_rgb=WHITE, text_rgb=INK, min_width=1.12)
    add_tag(slide, "ideal para clase", 8.5, 5.06, fill_rgb=WHITE, text_rgb=INK, min_width=1.2)
    add_footer(slide, "Esta es la bisagra conceptual para pasar a la clase practica")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, WHITE)
    add_decor(slide, GOLD)
    add_title(
        slide,
        "Guion visual de la hora",
        "En vez de recorrer definiciones en fila, usar una secuencia corta de escenas.",
        tag="RITMO",
    )
    add_timeline(
        slide,
        [
            ("0-5", "intuicion", CLAY),
            ("5-15", "tokens", BLUE),
            ("15-25", "contexto", GOLD),
            ("25-35", "prob + temp", TEAL),
            ("35-45", "transformer", CLAY),
            ("45-60", "local", GREEN),
        ],
        1.1,
        2.15,
        10.85,
    )
    add_card(slide, 0.95, 4.0, 3.55, 1.6, fill_rgb=CLAY_SOFT, line_rgb=CLAY)
    add_text(slide, "Primera parada", 1.2, 4.22, 1.3, 0.22, font_size=15, bold=True, color=CLAY)
    add_text(slide, "Hacer que repitan la idea de siguiente token antes de definir nada.", 1.2, 4.56, 2.9, 0.42, font_size=12.5)
    add_card(slide, 4.9, 4.0, 3.55, 1.6, fill_rgb=BLUE_SOFT, line_rgb=BLUE)
    add_text(slide, "Segunda parada", 5.15, 4.22, 1.4, 0.22, font_size=15, bold=True, color=BLUE)
    add_text(slide, "Mostrar una sola frase partida de tres maneras y preguntar que cambia.", 5.15, 4.56, 2.9, 0.42, font_size=12.5)
    add_card(slide, 8.85, 4.0, 3.55, 1.6, fill_rgb=TEAL_SOFT, line_rgb=TEAL)
    add_text(slide, "Ultima parada", 9.1, 4.22, 1.3, 0.22, font_size=15, bold=True, color=TEAL)
    add_text(slide, "Cerrar enlazando cuantizacion con la practica local de la clase siguiente.", 9.1, 4.56, 2.9, 0.42, font_size=12.5)
    add_footer(slide, "La clase mejora mucho si cada bloque deja una imagen facil de recordar")

    slide = prs.slides.add_slide(layout)
    set_bg(slide, PAPER)
    add_decor(slide, CLAY)
    add_title(
        slide,
        "Chequeo de salida",
        "Si al final dicen estas tres frases con sus palabras, la base quedo bien plantada.",
        tag="CIERRE",
    )
    takeaways = [
        (0.85, CLAY_SOFT, CLAY, "Un LLM predice el siguiente token, no responde por magia."),
        (4.45, BLUE_SOFT, BLUE, "Token no siempre es palabra: depende de como partimos el texto."),
        (8.05, TEAL_SOFT, TEAL, "llama.cpp importa porque vuelve local y usable lo que hoy vimos en teoria."),
    ]
    for left, fill_rgb, tone, sentence in takeaways:
        add_card(slide, left, 2.15, 3.2, 2.45, fill_rgb=fill_rgb, line_rgb=tone)
        add_text(slide, sentence, left + 0.22, 2.56, 2.75, 1.25, font_size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 0.96, 5.25, 11.3, 0.82, fill_rgb=WHITE, line_rgb=WHITE)
    add_text(
        slide,
        "Confusiones para cortar en el momento: piensa como humano | token = palabra | temperatura = mas inteligencia",
        1.18,
        5.53,
        10.8,
        0.22,
        font_size=14,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Fin del deck | menos enumeracion, mas referencias para senalar en clase")

    prs.save(OUTPUT)
    print(f"PPT generado en: {OUTPUT}")


if __name__ == "__main__":
    build_presentation()